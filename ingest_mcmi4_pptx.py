"""
Ingest MCMI-IV PowerPoint lecture content into Neon Postgres (pgvector) for RAG.

Extracts slide text, chunks it, generates embeddings with sentence-transformers,
and APPENDS to the existing mcmi4.chunks table (does not clear existing data).

Usage:
    python ingest_mcmi4_pptx.py <pptx_file>
"""

import sys
import os
import re
from pathlib import Path

os.environ["CUDA_VISIBLE_DEVICES"] = ""

from dotenv import load_dotenv
load_dotenv(Path(__file__).parent / ".env")

import psycopg2
from pptx import Presentation
from sentence_transformers import SentenceTransformer

DATABASE_URL = os.getenv("DATABASE_URL")
EMBEDDING_MODEL = "sentence-transformers/all-MiniLM-L6-v2"
CHUNK_SIZE = 500
CHUNK_OVERLAP = 100
SCHEMA = "mcmi4"


def extract_slides(pptx_path):
    """Extract text from each slide, return list of (slide_num, text) tuples."""
    prs = Presentation(pptx_path)
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                texts.append(shape.text.strip())
        # Deduplicate (title often repeats in body)
        seen = set()
        unique = []
        for t in texts:
            if t not in seen:
                seen.add(t)
                unique.append(t)
        slide_text = '\n'.join(unique).strip()
        if slide_text:
            slides.append((i, slide_text))
    return slides


def chunk_text(text, chunk_size=CHUNK_SIZE, overlap=CHUNK_OVERLAP):
    """Split text into overlapping chunks by approximate token count."""
    words = text.split()
    chunks = []
    start = 0
    while start < len(words):
        end = start + chunk_size
        chunk = ' '.join(words[start:end]).strip()
        if len(chunk) > 50:
            chunks.append(chunk)
        start += chunk_size - overlap
    return chunks


def main():
    if len(sys.argv) < 2:
        print("Usage: python ingest_mcmi4_pptx.py <pptx_file>")
        sys.exit(1)

    pptx_path = Path(sys.argv[1])
    if not pptx_path.exists():
        print(f"File not found: {pptx_path}")
        sys.exit(1)

    if not DATABASE_URL:
        print("DATABASE_URL not set in .env")
        sys.exit(1)

    # Extract slides
    print(f"Extracting slides from {pptx_path.name}...")
    slides = extract_slides(pptx_path)
    print(f"  {len(slides)} slides with content")

    # Combine all slide text with slide markers
    full_text = ""
    for slide_num, text in slides:
        full_text += f"\n\n[Slide {slide_num}]\n{text}\n"

    print(f"  Total characters: {len(full_text):,}")
    print(f"  Total words: {len(full_text.split()):,}")

    # Chunk
    chunks = chunk_text(full_text)
    print(f"  Generated {len(chunks)} chunks")

    # Connect to Neon
    print("Connecting to Neon...")
    conn = psycopg2.connect(DATABASE_URL)

    # Get current max chunk_index to append after existing data
    with conn.cursor() as cur:
        cur.execute(f"SELECT COALESCE(MAX(chunk_index), -1) FROM {SCHEMA}.chunks;")
        max_idx = cur.fetchone()[0]
        cur.execute(f"SELECT COUNT(*) FROM {SCHEMA}.chunks;")
        existing = cur.fetchone()[0]
    print(f"  Existing chunks: {existing} (max index: {max_idx})")

    # Generate embeddings
    print(f"Loading embedding model ({EMBEDDING_MODEL})...")
    model = SentenceTransformer(EMBEDDING_MODEL, device='cpu')

    print(f"Generating embeddings for {len(chunks)} chunks...")
    embeddings = model.encode(chunks, show_progress_bar=True, convert_to_numpy=True)

    # Insert into Neon (append with source_type='pptx_lecture')
    print(f"Inserting into {SCHEMA}.chunks...")
    start_idx = max_idx + 1
    with conn.cursor() as cur:
        for i, (chunk, embedding) in enumerate(zip(chunks, embeddings)):
            embedding_str = "[" + ",".join(str(float(x)) for x in embedding) + "]"
            cur.execute(
                f"INSERT INTO {SCHEMA}.chunks (content, chunk_index, source_type, embedding) "
                f"VALUES (%s, %s, %s, %s::vector)",
                (chunk, start_idx + i, 'pptx_lecture', embedding_str)
            )

    conn.commit()
    conn.close()

    print(f"\nDone! Appended {len(chunks)} chunks (source_type='pptx_lecture') to {SCHEMA}.chunks")
    print(f"Total chunks now: {existing + len(chunks)}")


if __name__ == '__main__':
    main()
